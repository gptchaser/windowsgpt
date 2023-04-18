from keys import OPENAI_KEY
import platform
import os

USERNAME = os.getlogin()
CURDIR = os.getcwd()
OPERATING_SYSTEM = platform.system()
PYTHON_VERSION = platform.python_version()
# in need of good prompt engineering
ENDOFTEXT = "<|ENDOFTEXT|>"
CODE_SYSTEM_CALIBRATION_MESSAGE = ENDOFTEXT+f"""您是 PythonGPT,由 OpenAI 训练的大型语言模型。请编写完整的 {OPERATING_SYSTEM} Python {PYTHON_VERSION} 代码，以便用户（用户名：{USERNAME}）可以运行它来解决其问题。在 ``` 区块中返回代码，不提供任何解释。不返回任何不是 Python 代码的文本。
Import 所有所需的requirements。"""
DEBUG_SYSTEM_CALIBRATION_MESSAGE = ENDOFTEXT+f"""你是PythonGPT,一个由OpenAI训练的大型语言模型。请编写完整的 {OPERATING_SYSTEM} Python {PYTHON_VERSION} 代码,以便用户可以运行它来解决他们的问题。例如,如果错误提示是“没有这样的文件或目录”或者“No such file or directory”, 那么您需要下载必要的文件或创建目录。请用简单明了的语言解释你的理由，然后提供更正后的代码。将整个代码都放在一个```块中。"""
INSTALL_SYSTEM_CALIBRATION_MESSAGE = ENDOFTEXT+"""你是PipGPT,一个由OpenAI训练的大型语言模型。请返回pip install命令以解决用户的问题。只返回命令,不要返回其他内容。"""
INSTALL_USER_MESSAGE = lambda package: f"""请编写 {OPERATING_SYSTEM} Python {PYTHON_VERSION} 的 pip 命令，以便我可以安装 {package}。请不要解释,只返回单个pip命令以进行安装。"""
LLM_SYSTEM_CALIBRATION_MESSAGE = """您是一位乐于助人的助手。请回应用户的目标。"""
CONGNITIVE_USER_MESSAGE = """使用提示工程和openai.Completion.create函数来帮助实现这个目标。使用text-davinci-003模型。不要忘记为语言模型设计提示,以便它返回相关的答案。"""
def USER_MESSAGE(goal, current_dir): return f"""（用户：{USERNAME}）
（目录：{current_dir}）
编写 {OPERATING_SYSTEM} python {PYTHON_VERSION} 代码，以便我可以运行代码实现我的目标。不要解释任何内容。仅返回代码。我的目标是：[{goal}]。不要忘记打印最终结果。"""
def DEBUG_MESSAGE(code, error):
    return f"""```python
{code}
```
上述代码返回错误 "{error}"。请简要解释为什么会出现这个错误，然后编写更正后的完整代码。""" 
CODE_USER_CALIBRATION_MESSAGE = """获取有关爱丁顿亮度的信息，然后制作一个关于它的 PowerPoint。"""
CODE_ASSISTANT_CALIBRATION_MESSAGE = """```python
import wikipedia
import pptx
import openai
openai.api_key = "your_openai_api_key_here"

# Set the language to English
wikipedia.set_lang("en")

# Get the page object for Artificial Neural Networks (we never want auto_suggest)
ann_page = wikipedia.page("Eddington Luminosity", auto_suggest=False)

# Create a new PowerPoint presentation
prs = pptx.Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Eddington Luminosity"

# Add a slide for each section of the Wikipedia page
for section in ann_page.sections:
    # Skip the first section ("Overview")
    if section.title == "Overview":
        continue
    # Add a new slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    # Set the title of the slide to the section title
    slide.shapes.title.text = section
    # Use language model to make bullet points
    bullet_points = slide.shapes.placeholders[1]
    section_text = ann_page.section(section)
    prompt = f"Information is given in the following square brackets: [{section_text}]. Please respond with very brief presentation slide bullet points for it, separated by a ;."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        temperature=0.5,
        max_tokens=512,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    print(response.choices[0].text)
    for point in response_text.split(';'):
        # Add the bullet point to the slide
        bullet_item = bullet_points.text_frame.add_paragraph()
        bullet_item.text = point

# Save the PowerPoint presentation
prs.save("Eddington_Luminosity.pptx")

# Print to confirm goal has been completed
print("PowerPoint presentation Eddington_Luminosity.pptx created.")```"""
CODE_USER_CALIBRATION_MESSAGE_UNSPLASH_EXAMPLE = """make my wallpaper a galaxy"""
CODE_ASSISTANT_CALIBRATION_MESSAGE_UNSPLASH_EXAMPLE = """```python
import requests
import ctypes
import os
url = "https://api.unsplash.com/search/photos"
params = {
    "query": "galaxy",    # search for "galaxy"
    "orientation": "landscape",   # limit results to landscape orientation
    "client_id": "your_unsplash_access_key_here"   # Unsplash access key
}
response = requests.get(url, params=params)
# Get the URL of the first image in the results
image_url = response.json()["results"][0]["urls"]["regular"]
# Download the image and save it to a file
response = requests.get(image_url)
with open("galaxy.jpg", "wb") as f:
    f.write(response.content)
# Change it to a galaxy
ctypes.windll.user32.SystemParametersInfoW(20, 0, os.path.abspath("galaxy.jpg"), 3)
# Print to confirm goal has been completed
print("Wallpaper changed to a galaxy.")```"""
CONSOLE_OUTPUT_CALIBRATION_MESSAGE = """PowerPoint presentation Eddington_Luminosity.pptx created."""
CONSOLE_OUTPUT_CALIBRATION_MESSAGE_UNSPLASH_EXAMPLE = """Wallpaper changed to a galaxy."""
